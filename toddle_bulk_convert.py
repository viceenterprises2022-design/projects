"""Convert all downloaded subject files to text for NotebookLM."""
import sys
from pathlib import Path
from docx import Document
from pptx import Presentation
from pdfminer.high_level import extract_text as pdf_extract

SRC = Path("output/downloads")
OUT = Path("output/text")

def extract_docx(path):
    return "\n".join(p.text for p in Document(path).paragraphs if p.text.strip())

def extract_pptx(path):
    texts = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t: texts.append(t)
    return "\n".join(texts)

def extract_pdf(path):
    return pdf_extract(path)


def convert_subject(subj):
    sdir = SRC / subj
    if not sdir.exists():
        print(f"SKIP {subj}: no download dir")
        return 0

    odir = OUT / subj
    odir.mkdir(parents=True, exist_ok=True)
    count = 0

    for fpath in sorted(sdir.rglob("*")):
        if not fpath.is_file():
            continue

        ext = fpath.suffix.lower()
        if ext in (".jpg", ".jpeg", ".png", ".gif"):
            out_path = odir / fpath.relative_to(sdir).with_suffix(".md")
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(f"*Image: {fpath.name}*\n")
            count += 1
            continue

        if ext == ".docx":
            kind = "docx"
        elif ext == ".pptx":
            kind = "pptx"
        elif ext == ".pdf":
            kind = "pdf"
        else:
            continue

        try:
            if kind == "docx": text = extract_docx(fpath)
            elif kind == "pptx": text = extract_pptx(fpath)
            else: text = extract_pdf(fpath)
        except Exception as e:
            print(f"  ERR {fpath.relative_to(SRC)}: {e}")
            continue

        rel = fpath.relative_to(sdir)
        out_path = odir / rel.with_suffix(".md")
        out_path.parent.mkdir(parents=True, exist_ok=True)

        lines = [l.strip() for l in text.split("\n") if l.strip()]
        clean = "\n".join(lines)
        front = f"---\nsource: {fpath.name}\ntype: {kind}\n---\n\n"
        out_path.write_text(front + clean)
        print(f"  OK  {rel} ({len(clean)}b)")
        count += 1

    return count


if __name__ == "__main__":
    subjects = ["Chemistry", "Mathematics", "English", "Biology", "History", "Geography", "Spanish", "Design", "Visual Arts"]
    grand = 0

    for subj in subjects:
        print(f"\n{'='*50}\n{subj}")
        c = convert_subject(subj)
        print(f"  {c} files converted")
        grand += c

    print(f"\nTotal: {grand} files converted")
