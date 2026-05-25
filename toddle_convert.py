"""Convert downloaded Toddle files to text for NotebookLM ingestion."""
import json, re
from pathlib import Path
from docx import Document
from pptx import Presentation
from pdfminer.high_level import extract_text as pdf_extract

SRC = Path("output/downloads/physics")
OUT = Path("output/text/physics")

EXT_MAP = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".txt": "txt",
    ".md": "txt",
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".gif": "image",
}

def extract_docx(path):
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)

def extract_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return "\n".join(texts)

def extract_pdf(path):
    return pdf_extract(path)

def is_binary_image(path):
    ext = path.suffix.lower()
    return ext in (".jpg", ".jpeg", ".png", ".gif")


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    total_md = 0
    total_bytes = 0

    for fpath in sorted(SRC.rglob("*")):
        if not fpath.is_file():
            continue

        ext = fpath.suffix.lower()
        kind = EXT_MAP.get(ext)

        if kind == "image":
            rel = fpath.relative_to(SRC)
            out_path = OUT / rel.with_suffix(".md")
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(f"*Image file: {fpath.name}*\n")
            total_md += 1
            continue

        if kind not in ("docx", "pptx", "pdf"):
            continue

        try:
            if kind == "docx":
                text = extract_docx(fpath)
            elif kind == "pptx":
                text = extract_pptx(fpath)
            elif kind == "pdf":
                text = extract_pdf(fpath)
            else:
                continue
        except Exception as e:
            print(f"ERR extracting {fpath.relative_to(SRC)}: {e}")
            continue

        rel = fpath.relative_to(SRC)
        out_path = OUT / rel.with_suffix(".md")
        out_path.parent.mkdir(parents=True, exist_ok=True)

        lines = text.split("\n")
        clean = [l.strip() for l in lines if l.strip()]
        clean_text = "\n".join(clean)

        front_matter = (
            f"---\n"
            f"source: {fpath.name}\n"
            f"type: {kind}\n"
            f"---\n\n"
        )

        out_path.write_text(front_matter + clean_text)

        size = len(clean_text.encode("utf-8"))
        pages = max(1, size // 3000)
        total_md += 1
        total_bytes += size

        if size < 100:
            print(f"  TINY {fpath.name} ({size}b)")
        else:
            print(f"  OK   {fpath.name} ({size}b, ~{pages}p)")

    print(f"\nDone: {total_md} files converted, {total_bytes/1024:.0f} KB total")

if __name__ == "__main__":
    main()
